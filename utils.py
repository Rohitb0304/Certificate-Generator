import os
import tempfile
import subprocess
from pptx import Presentation
from pptx.util import Pt  # For handling font size
from pptx.enum.text import PP_ALIGN  # For text alignment

def generate_certificate(row, pptx_template_path, output_folder):
    try:
        full_name = row['Full_Name'].title()
        reference_number = row.get('Reference_Number', 'N/A').strip()

        # Load the PowerPoint template
        prs = Presentation(pptx_template_path)

        # Replace placeholders in the slide
        slide = prs.slides[0]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Check for <<FULL NAME>> tag and replace it with the student's name
                    if '<<FULL NAME>>' in run.text:
                        original_text = run.text
                        run.text = original_text.replace('<<FULL NAME>>', full_name)

                        # Keep the font size, bold, italics, and other styles from the template
                        run.font.size = run.font.size or Pt(24)  # Set a default size if not present
                        run.font.bold = run.font.bold
                        run.font.italic = run.font.italic
                        run.font.underline = run.font.underline

                    # Check for <<REFERENCE NO>> tag and replace it with the reference number
                    if '<<REFERENCE NO>>' in run.text:
                        original_text = run.text
                        run.text = original_text.replace('<<REFERENCE NO>>', str(reference_number))

        # Save the modified certificate as a new PPTX file
        temp_pptx_path = os.path.join(output_folder, f"{full_name}_{reference_number}.pptx")
        prs.save(temp_pptx_path)

        # Convert to PDF (optional)
        convert_to_pdf(temp_pptx_path, output_folder)

    except Exception as e:
        print(f"Error generating certificate for {full_name}: {e}")


def convert_to_pdf(pptx_file, output_folder):
    try:
        subprocess.run(['/usr/bin/libreoffice', '--headless', '--convert-to', 'pdf', pptx_file, '--outdir', output_folder], check=True)
        os.remove(pptx_file)  # Optionally delete the .pptx file
    except subprocess.CalledProcessError as e:
        print(f"Failed to convert {pptx_file} to PDF: {e}")
