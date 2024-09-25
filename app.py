from flask import Flask, request, render_template, send_file, redirect, url_for, flash, session, jsonify
import pandas as pd
from pptx import Presentation
import os
import zipfile
import shutil
from utils import generate_certificate
import time

app = Flask(__name__)
app.secret_key = 'your_secret_key'

# File paths
UPLOAD_FOLDER = 'uploads/'
CERTIFICATES_FOLDER = 'certificates/'
ZIP_FOLDER = 'zips/'

# Create folders if they don't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(CERTIFICATES_FOLDER, exist_ok=True)
os.makedirs(ZIP_FOLDER, exist_ok=True)

# Global progress tracking
progress_status = {
    "total": 0,
    "completed": 0,
    "current_name": ""
}

# Store processed rows to prevent infinite loop or repeat
processed_rows = set()

def check_pptx_tags(pptx_path):
    """Check if the PPTX file contains the tag <<FULL NAME>>"""
    prs = Presentation(pptx_path)
    tags_found = any('<<FULL NAME>>' in shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame)
    return tags_found

def apply_name_with_template_style(shape, name):
    """Replaces <<FULL NAME>> with the actual name, keeping the original style (font, size, color)"""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if '<<FULL NAME>>' in run.text:
                    # Store original font properties
                    original_font = run.font
                    font_size = original_font.size
                    font_color = original_font.color.rgb
                    bold = original_font.bold
                    italic = original_font.italic

                    # Replace <<FULL NAME>> with the actual student name
                    run.text = name
                    
                    # Apply the original formatting
                    run.font.size = font_size
                    run.font.color.rgb = font_color
                    run.font.bold = bold
                    run.font.italic = italic

def generate_certificate_with_fonts(row, pptx_path, output_folder):
    """Generates a certificate for a student and preserves the font from PPTX."""
    prs = Presentation(pptx_path)
    student_name = row['Full_Name']

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                apply_name_with_template_style(shape, student_name)

    output_file = os.path.join(output_folder, f"{student_name}.pptx")
    prs.save(output_file)

@app.route('/', methods=['GET', 'POST'])
def index():
    global progress_status
    global processed_rows

    if request.method == 'POST':
        if 'pptx_file' not in request.files or 'csv_file' not in request.files:
            flash('Please upload both PPTX and CSV files!')
            return redirect(request.url)

        pptx_file = request.files['pptx_file']
        csv_file = request.files['csv_file']

        # Save uploaded files
        pptx_path = os.path.join(UPLOAD_FOLDER, pptx_file.filename)
        csv_path = os.path.join(UPLOAD_FOLDER, csv_file.filename)

        pptx_file.save(pptx_path)
        csv_file.save(csv_path)

        if not check_pptx_tags(pptx_path):
            flash('Tag <<FULL NAME>> not found in PPTX template.')
            return redirect(request.url)

        flash('Tag <<FULL NAME>> found. Now fetching CSV data...')

        # Load CSV
        student_data = pd.read_csv(csv_path)
        if len(student_data) > 80:
            flash('The limit is 80 records. Please upload a smaller file.')
            return redirect(request.url)

        # Initialize progress status
        progress_status = {
            "total": len(student_data),
            "completed": 0,
            "current_name": ""
        }

        # Process each student and update progress
        for idx, row in student_data.iterrows():
            student_name = row['Full_Name']
            progress_status['current_name'] = student_name  # Update the current student's name
            
            if idx in processed_rows:
                continue  # Skip already processed rows

            # Generate certificate and track the processed row
            generate_certificate_with_fonts(row, pptx_path, CERTIFICATES_FOLDER)
            processed_rows.add(idx)  # Mark this row as processed
            
            progress_status['completed'] += 1  # Update the count of completed certificates
            time.sleep(1)  # Simulate delay for certificate generation

        flash('Certificates generated successfully! Preparing download...')

        # Create a ZIP of the certificates
        zip_filename = os.path.join(ZIP_FOLDER, 'certificates.zip')
        with zipfile.ZipFile(zip_filename, 'w') as zipf:
            for root, dirs, files in os.walk(CERTIFICATES_FOLDER):
                for file in files:
                    zipf.write(os.path.join(root, file), os.path.relpath(os.path.join(root, file), CERTIFICATES_FOLDER))

        # Clear old certificates
        shutil.rmtree(CERTIFICATES_FOLDER)
        os.makedirs(CERTIFICATES_FOLDER, exist_ok=True)

        return send_file(zip_filename, as_attachment=True, download_name='certificates.zip')

    return render_template('index.html')

@app.route('/progress', methods=['GET'])
def progress():
    """Return the current progress of certificate generation."""
    return jsonify(progress_status)

if __name__ == "__main__":
    app.run(host='0.0.0.0', port=5000, debug=True)
